import os
import time
import json
import logging
import mimetypes
from typing import List, Dict, Optional, Tuple
import io

import requests
from dotenv import load_dotenv
from pptx import Presentation
from openai import OpenAI

# 导入必要SDK并检查
try:
    import cloudinary
    import cloudinary.uploader
    from cloudinary.utils import cloudinary_url
except ImportError:
    raise ImportError("请安装Cloudinary SDK：pip install cloudinary")

try:
    import asposeslidescloud
    from asposeslidescloud.configuration import Configuration
    from asposeslidescloud.apis.slides_api import SlidesApi
except ImportError:
    raise ImportError("请安装Aspose SDK：pip install asposeslidescloud")


# 初始化环境与日志
load_dotenv()
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class CloudinaryStorage:
    """简化的Cloudinary存储管理（仅保留核心文件操作）"""
    def __init__(self):
        # 读取并验证Cloudinary配置
        self.cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME")
        self.api_key = os.getenv("CLOUDINARY_API_KEY")
        self.api_secret = os.getenv("CLOUDINARY_API_SECRET")
        
        if not all([self.cloud_name, self.api_key, self.api_secret]):
            raise ValueError("缺少Cloudinary环境变量，请检查.env文件")

        # 初始化Cloudinary客户端
        cloudinary.config(
            cloud_name=self.cloud_name,
            api_key=self.api_key,
            api_secret=self.api_secret,
            secure=True
        )
        logging.info("Cloudinary存储初始化完成")

    def get_file_bytes(self, public_id: str) -> Optional[bytes]:
        """从Cloudinary获取文件字节流（支持PPTX等非图片文件）"""
        try:
            file_url, _ = cloudinary_url(public_id, resource_type="raw")
            response = requests.get(file_url, timeout=60)
            response.raise_for_status()  # 触发HTTP错误（如404、500）
            return response.content
        except Exception as e:
            logging.error(f"获取文件[{public_id}]失败：{e}")
            return None

    def upload_file(self, file_bytes: bytes, public_id: str) -> Tuple[bool, str]:
        """上传文件到Cloudinary，返回（成功状态，public_id/错误信息）"""
        try:
            result = cloudinary.uploader.upload(
                io.BytesIO(file_bytes),
                public_id=public_id,
                resource_type="raw",  # 非媒体文件用raw
                overwrite=True
            )
            return True, result["public_id"]
        except Exception as e:
            error_msg = f"上传文件[{public_id}]失败：{e}"
            logging.error(error_msg)
            return False, error_msg


class PPTXToHeyGenVideo:
    """简化的PPTX转HeyGen视频工具（合并重复逻辑，扁平流程）"""
    def __init__(self, storage_manager: CloudinaryStorage):
        # 1. 绑定存储管理器
        self.storage = storage_manager
        
        # 2. 读取并验证所有API密钥
        self.heygen_key = os.getenv("HEYGEN_API_KEY")
        self.openai_key = os.getenv("OPENAI_API_KEY")
        self.aspose_id = os.getenv("ASPOSE_CLIENT_ID")
        self.aspose_secret = os.getenv("ASPOSE_CLIENT_SECRET")
        self._validate_config()

        # 3. 初始化API客户端
        self.openai_client = OpenAI(api_key=self.openai_key)
        self.heygen_headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Api-Key": self.heygen_key
        }
        # 初始化Aspose（用于PPTX转图片）
        aspose_config = Configuration()
        aspose_config.app_sid = self.aspose_id
        aspose_config.app_key = self.aspose_secret
        self.aspose_slides_api = SlidesApi(aspose_config)

        # 4. 视频基础配置（可根据需求调整）
        self.avatar_id = os.getenv("HEYGEN_AVATAR_ID")
        self.voice_id = os.getenv("HEYGEN_VOICE_ID")
        self.video_size = {"width": 1280, "height": 720}
        self.retry_times = 3  # 网络请求重试次数
        self.poll_interval = 10  # 视频生成状态查询间隔（秒）
        self.slide_asset_ids = []  # 存储HeyGen中的幻灯片图片ID

    def _validate_config(self):
        """验证所有必要配置，缺失则报错"""
        required = [
            ("HEYGEN_API_KEY", self.heygen_key),
            ("OPENAI_API_KEY", self.openai_key),
            ("ASPOSE_CLIENT_ID", self.aspose_id),
            ("ASPOSE_CLIENT_SECRET", self.aspose_secret),
            ("HEYGEN_AVATAR_ID", self.avatar_id),
            ("HEYGEN_VOICE_ID", self.voice_id)
        ]
        for name, value in required:
            if not value:
                raise ValueError(f"缺少环境变量：{name}，请检查.env文件")

    def _request_with_retry(self, method: str, url: str, **kwargs) -> requests.Response:
        """通用网络请求重试逻辑（合并原_post/_get_with_retry）"""
        for attempt in range(self.retry_times):
            try:
                response = requests.request(method, url, timeout=90, **kwargs)
                response.raise_for_status()
                return response
            except requests.RequestException as e:
                logging.warning(f"请求[{method} {url}]第{attempt+1}次失败：{e}")
                if attempt == self.retry_times - 1:
                    raise  # 最后一次失败则抛出异常
                time.sleep(2 ** attempt)  # 指数退避等待
        raise RuntimeError("请求重试次数耗尽")

    def _upload_to_heygen(self, file_bytes: bytes, file_name: str) -> str:
        """上传文件（如幻灯片图片）到HeyGen，返回资产ID"""
        # 自动识别文件类型（默认png）
        mime_type = mimetypes.guess_type(file_name)[0] or "image/png"
        upload_headers = {"X-Api-Key": self.heygen_key, "Content-Type": mime_type}

        # 用通用重试逻辑上传
        response = self._request_with_retry(
            method="POST",
            url="https://upload.heygen.com/v1/asset",
            headers=upload_headers,
            data=io.BytesIO(file_bytes)
        )
        asset_id = response.json()["data"].get("id")
        if not asset_id:
            raise RuntimeError(f"HeyGen上传[{file_name}]未返回资产ID")
        
        logging.info(f"HeyGen上传成功：{file_name} → 资产ID：{asset_id}")
        return asset_id

    def _pptx_to_heygen_images(self, pptx_bytes: bytes, slide_count: int):
        """将PPTX幻灯片转为图片并上传到HeyGen，更新slide_asset_ids"""
        logging.info(f"开始转换{slide_count}张幻灯片为图片")
        for slide_idx in range(1, slide_count + 1):  # Aspose幻灯片索引从1开始
            temp_file_path = None
            try:
                # 1. 用Aspose将单张幻灯片转为PNG（返回临时文件路径）
                temp_file_path = self.aspose_slides_api.download_slide_online(
                    document=io.BytesIO(pptx_bytes),
                    slide_index=slide_idx,
                    format="PNG"
                )
                # 2. 读取临时文件字节流
                with open(temp_file_path, "rb") as f:
                    slide_bytes = f.read()
                if not slide_bytes:
                    raise ValueError(f"第{slide_idx}张幻灯片转换后为空")
                # 3. 上传到HeyGen并记录资产ID
                asset_id = self._upload_to_heygen(slide_bytes, f"slide_{slide_idx}.png")
                self.slide_asset_ids.append(asset_id)
            except Exception as e:
                logging.error(f"处理第{slide_idx}张幻灯片失败：{e}")
            finally:
                # 清理Aspose生成的临时文件
                if temp_file_path and os.path.exists(temp_file_path):
                    os.remove(temp_file_path)

    def _generate_speaker_notes(self, slide_texts: List[str]) -> List[str]:
        """用OpenAI生成每张幻灯片的讲稿（虚拟教师风格）"""
        logging.info(f"为{len(slide_texts)}张幻灯片生成讲稿")
        notes = []
        for idx, text in enumerate(slide_texts, 1):
            if not text.strip():
                notes.append(f"这是第{idx}张幻灯片，主要展示图片内容。")
                continue
            try:
                # 调用OpenAI API生成讲稿
                response = self.openai_client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "你是虚拟教师，用简洁易懂的语言讲解幻灯片内容，语气专业亲切。"},
                        {"role": "user", "content": f"根据以下幻灯片文本生成讲稿：\n{text}"}
                    ],
                    temperature=0.7,
                    max_tokens=300
                )
                notes.append(response.choices[0].message.content.strip())
            except Exception as e:
                logging.error(f"生成第{idx}张幻灯片讲稿失败：{e}")
                notes.append(f"第{idx}张幻灯片内容：{text[:50]}...")  # 降级处理
        return notes

    def _create_heygen_video(self, speaker_notes: List[str], video_title: str) -> str:
        """创建HeyGen多场景视频，返回视频ID"""
        # 1. 构建每个场景的配置（虚拟人+背景图+讲稿）
        video_scenes = []
        for idx, note in enumerate(speaker_notes):
            scene = {
                "character": {
                    "type": "talking_photo",
                    "talking_photo_id": self.avatar_id,
                    "scale": 0.33,
                    "offset": {"x": 0.42, "y": 0.42}
                },
                "voice": {"type": "text", "input_text": note, "voice_id": self.voice_id},
                "background": {"type": "color", "value": "#FFFFFF"}  # 默认白色背景
            }
            # 若有对应幻灯片图片，替换背景
            if idx < len(self.slide_asset_ids):
                scene["background"] = {"type": "image", "image_asset_id": self.slide_asset_ids[idx]}
            video_scenes.append(scene)

        # 2. 调用HeyGen API创建视频
        payload = {
            "video_inputs": video_scenes,
            "dimension": self.video_size,
            "title": video_title
        }
        response = self._request_with_retry(
            method="POST",
            url="https://api.heygen.com/v2/video/generate",
            headers=self.heygen_headers,
            json=payload
        )
        video_id = response.json()["data"].get("video_id")
        if not video_id:
            raise RuntimeError("HeyGen未返回视频ID，创建失败")
        logging.info(f"HeyGen视频创建成功，视频ID：{video_id}")
        return video_id

    def _wait_for_video(self, video_id: str) -> str:
        """等待视频生成完成，返回最终视频URL"""
        logging.info(f"开始等待视频[{video_id}]生成，每{self.poll_interval}秒查询一次状态")
        while True:
            # 查询视频状态
            response = self._request_with_retry(
                method="GET",
                url=f"https://api.heygen.com/v1/video_status.get?video_id={video_id}",
                headers=self.heygen_headers
            )
            data = response.json()["data"]
            status = data["status"]

            if status in ("completed", "success"):
                return data["video_url"]  # 生成成功，返回URL
            elif status in ("failed", "error"):
                raise RuntimeError(f"视频生成失败：{data.get('error', '未知错误')}")
            else:
                logging.info(f"视频状态：{status}，继续等待...")
                time.sleep(self.poll_interval)

    def convert(self, pptx_public_id: str, video_title: Optional[str] = None, max_slides: Optional[int] = None) -> Dict:
        """核心转换流程：Cloudinary获取PPTX → 转图片 → 生成讲稿 → 生成视频 → 返回结果"""
        # 1. 从Cloudinary获取PPTX文件
        logging.info(f"开始处理PPTX：{pptx_public_id}")
        pptx_bytes = self.storage.get_file_bytes(pptx_public_id)
        if not pptx_bytes:
            raise FileNotFoundError(f"Cloudinary中未找到PPTX：{pptx_public_id}")

        # 2. 解析PPTX，提取幻灯片文本（限制最大处理数量）
        prs = Presentation(io.BytesIO(pptx_bytes))
        slides = list(prs.slides)
        if not slides:
            raise ValueError("PPTX文件中没有幻灯片")
        if max_slides:
            slides = slides[:max_slides]  # 限制处理的最大幻灯片数
        # 提取每张幻灯片的文本内容
        slide_texts = [
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            for slide in slides
        ]

        # 3. 幻灯片转图片并上传到HeyGen
        self._pptx_to_heygen_images(pptx_bytes, len(slides))
        if not self.slide_asset_ids:
            raise RuntimeError("没有成功转换任何幻灯片为图片")

        # 4. 生成讲稿（确保讲稿数量与图片数量一致）
        speaker_notes = self._generate_speaker_notes(slide_texts[:len(self.slide_asset_ids)])

        # 5. 创建并等待视频生成
        final_title = video_title or f"PPTX转换视频_{pptx_public_id}"
        video_id = self._create_heygen_video(speaker_notes, final_title)
        video_url = self._wait_for_video(video_id)

        # 6. 返回最终结果
        return {
            "video_id": video_id,
            "video_url": video_url,
            "slides_processed": len(speaker_notes),
            "title": final_title
        }


# 端到端使用示例（直接运行脚本即可）
if __name__ == "__main__":
    try:
        # 1. 初始化存储
        cloud_storage = CloudinaryStorage()
        # 2. 替换为你在Cloudinary中的PPTX文件public_id
        target_pptx_id = "你的PPTX文件public_id"  # 例："presentations/my_class_ppt"
        # 3. 初始化转换器并执行转换
        converter = PPTXToHeyGenVideo(storage_manager=cloud_storage)
        result = converter.convert(
            pptx_public_id=target_pptx_id,
            video_title="我的PPT转视频示例",
            max_slides=10  # 可选：限制最大处理10张幻灯片
        )
        # 4. 打印结果
        print("\n=== 转换完成 ===")
        print(json.dumps(result, indent=2, ensure_ascii=False))
    except Exception as e:
        logging.error(f"转换过程出错：{e}", exc_info=True)
        print(f"执行失败：{e}")